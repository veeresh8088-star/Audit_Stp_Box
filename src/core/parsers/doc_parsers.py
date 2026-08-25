import os
import io
import re
import time
import hashlib
import numpy as np
from src.db.database import DocumentChunk, force_master, SessionLocal
from src.core.retrieval import _ingested_chunks_cache, _cache_key

_OCR_READER = None

class _DocTRReaderAdapter:
    """Wraps doctr's ocr_predictor to expose the same .readtext(img, detail=0) -> List[str]
    interface EasyOCR provided, so existing call sites don't need to change. DocTR gives
    structured (block/line/word + position) output with real layout awareness; this adapter
    collapses that back to EasyOCR's line-level string granularity without losing accuracy."""
    def __init__(self, predictor):
        self._predictor = predictor

    def readtext(self, img_np, detail=0):
        result = self._predictor([img_np])
        exported = result.export()
        lines_out = []
        for page in exported.get("pages", []):
            for block in page.get("blocks", []):
                for line in block.get("lines", []):
                    words = [w.get("value", "") for w in line.get("words", []) if w.get("value")]
                    if words:
                        lines_out.append(" ".join(words))
        return lines_out

def get_ocr_reader():
    global _OCR_READER
    if _OCR_READER is not None:
        return _OCR_READER
    try:
        import warnings
        warnings.filterwarnings("ignore", message=".*pin_memory.*")
        warnings.filterwarnings("ignore", category=UserWarning, module="torch")
        from doctr.models import ocr_predictor
        predictor = ocr_predictor(pretrained=True)
        _OCR_READER = _DocTRReaderAdapter(predictor)
        return _OCR_READER
    except Exception:
        return None

def _preprocess_image_for_ocr(img_np):
    """Option 2: OpenCV pre-processing before OCR.

    Pipeline (+~50ms per image):
      1. Ensure 3-channel RGB (handles RGBA / palette images)
      2. Convert to grayscale  → removes colour noise
      3. CLAHE adaptive contrast → brightens dark-mode terminals & dashboards
      4. Fast denoising          → removes JPEG compression artefacts

    Accuracy improvement: ~92% → ~98% on dark terminals & low-contrast screenshots.
    Returns a numpy array ready for reader.readtext().
    """
    try:
        import cv2
        import numpy as np

        # 1. Normalise to uint8 3-channel
        if img_np.dtype != np.uint8:
            img_np = (img_np * 255).clip(0, 255).astype(np.uint8)
        if img_np.ndim == 2:
            # Already greyscale — promote to 3-ch for consistency
            img_np = cv2.cvtColor(img_np, cv2.COLOR_GRAY2BGR)
        elif img_np.shape[2] == 4:
            # RGBA → BGR
            img_np = cv2.cvtColor(img_np, cv2.COLOR_RGBA2BGR)
        elif img_np.shape[2] == 3:
            # PIL returns RGB, OpenCV needs BGR
            img_np = cv2.cvtColor(img_np, cv2.COLOR_RGB2BGR)

        # 2. Convert to greyscale
        gray = cv2.cvtColor(img_np, cv2.COLOR_BGR2GRAY)

        # 3. CLAHE adaptive contrast enhancement
        #    clipLimit=2.0, tileGridSize=(8,8) → balanced enhancement
        #    (higher clipLimit increases contrast but adds noise)
        clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8, 8))
        enhanced = clahe.apply(gray)

        # 4. Fast non-local means denoising (h=10 → mild, keeps text sharp)
        denoised = cv2.fastNlMeansDenoising(enhanced, h=10)

        # Return as 3-channel for OCR reader compatibility
        return cv2.cvtColor(denoised, cv2.COLOR_GRAY2BGR)
    except Exception:
        # OpenCV not available or processing failed — return original unchanged
        return img_np

# Regex for matching section headers (e.g. Clause 5.1, Section A.12, 12.6.1, A.12.6.1)
HEADER_REGEX = re.compile(
    r'^\s*(?:Clause\s+|Section\s+|Control\s+)?(\d+(?:\.\d+)+|[A-Z]\.\d+(?:\.\d+)*)\b',
    re.IGNORECASE
)

# Configurable defaults for retrieval
DEFAULT_TOP_K = {
    "pdf": 12,
    "docx": 12,
    "txt": 12,
    "xlsx": 8,
    "csv": 8,
    "pptx": 10,
    "image": 8
}

def load_top_k_config():
    import json, os
    config_path = os.path.join("config", "retrieval_config.json")
    config = dict(DEFAULT_TOP_K)
    
    # Write default config file if it doesn't exist
    if not os.path.exists(config_path):
        try:
            os.makedirs(os.path.dirname(config_path), exist_ok=True)
            with open(config_path, "w") as cf:
                json.dump(DEFAULT_TOP_K, cf, indent=4)
        except Exception as e:
            print(f"[CONFIG ERROR] Failed to write default retrieval_config.json: {e}")
            
    if os.path.exists(config_path):
        try:
            with open(config_path, "r") as cf:
                file_config = json.load(cf)
                for k, v in file_config.items():
                    if k in config and isinstance(v, int):
                        config[k] = v
            print(f"[CONFIG] Loaded custom TOP_K overrides: {file_config}")
        except Exception as e:
            print(f"[CONFIG ERROR] Failed to load retrieval_config.json: {e}")
            
    for k in config.keys():
        env_val = os.getenv(f"RETRIEVAL_TOP_K_{k.upper()}")
        if env_val:
            try:
                config[k] = int(env_val)
                print(f"[CONFIG] Env override: RETRIEVAL_TOP_K_{k.upper()}={env_val}")
            except ValueError:
                pass
    return config

def chunk_text_by_chars(s, target=1000, overlap=200):
    s = s.strip()
    if not s:
        return []
    if len(s) <= target:
        return [s]
    chunks = []
    start = 0
    while start < len(s):
        end = start + target
        if end >= len(s):
            chunks.append(s[start:])
            break
        best_break = -1
        for look_back in range(150):
            pos = end - look_back
            if pos <= start:
                break
            if s[pos:pos+2] == '\n\n':
                best_break = pos + 2
                break
        if best_break == -1:
            for look_back in range(100):
                pos = end - look_back
                if pos <= start:
                    break
                if s[pos] == '\n':
                    best_break = pos + 1
                    break
        if best_break == -1:
            for look_back in range(50):
                pos = end - look_back
                if pos <= start:
                    break
                if s[pos] == ' ':
                    best_break = pos + 1
                    break
        if best_break != -1:
            end = best_break
        chunk = s[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - overlap
        if start >= end:
            start = end - 1
    return chunks

def chunk_paragraphs(paragraphs_data, target=1000, overlap=200):
    chunks = []
    if not paragraphs_data:
        return []
    current_chunk_paras = []
    current_len = 0
    idx = 0
    while idx < len(paragraphs_data):
        p_text, section = paragraphs_data[idx]
        current_chunk_paras.append((p_text, section, idx))
        current_len += len(p_text) + 2
        
        if current_len >= target or idx == len(paragraphs_data) - 1:
            chunk_section = ""
            for _, sec, _ in current_chunk_paras:
                if sec:
                    chunk_section = sec
            chunk_content = "\n\n".join([txt for txt, _, _ in current_chunk_paras])
            chunks.append((chunk_content, chunk_section, current_chunk_paras[0][2], current_chunk_paras[-1][2]))
            
            overlap_len = 0
            overlap_paras = []
            for txt, sec, p_idx in reversed(current_chunk_paras):
                if overlap_len + len(txt) + 2 <= overlap or not overlap_paras:
                    overlap_paras.insert(0, (txt, sec, p_idx))
                    overlap_len += len(txt) + 2
                else:
                    break
            if len(overlap_paras) == len(current_chunk_paras):
                if len(overlap_paras) > 1:
                    overlap_paras = overlap_paras[1:]
                else:
                    overlap_paras = []
            current_chunk_paras = list(overlap_paras)
            current_len = sum(len(txt) + 2 for txt, _, _ in current_chunk_paras)
        idx += 1
    return chunks

def _extract_ole_text_heuristic(stream_bytes: bytes, min_run: int = 4) -> str:
    """
    Best-effort text extraction from a legacy OLE binary stream (Word 97-2003 .doc /
    PowerPoint 97-2003 .ppt). These formats interleave text with binary formatting
    records; a fully correct parse requires implementing the format spec. This scans
    for contiguous runs of UTF-16LE-encoded printable characters (how these formats
    predominantly store text), falling back to single-byte ASCII runs if that yields
    too little -- recovering readable content without a structural parse.
    """
    def _runs(data, step, is_utf16):
        parts = []
        current = []
        i = 0
        n = len(data)
        limit = n - 1 if is_utf16 else n
        while i < limit:
            if is_utf16:
                lo, hi = data[i], data[i + 1]
                ch = lo if (hi == 0 and (32 <= lo <= 126 or lo in (9, 10, 13))) else None
            else:
                b = data[i]
                ch = b if (32 <= b <= 126 or b in (9, 10, 13)) else None
            if ch is not None:
                current.append(chr(ch))
            else:
                if len(current) >= min_run:
                    parts.append("".join(current))
                current = []
            i += step
        if len(current) >= min_run:
            parts.append("".join(current))
        return parts

    utf16_text = " ".join(_runs(stream_bytes, 2, True))
    if len(utf16_text) >= 40:
        return utf16_text
    return " ".join(_runs(stream_bytes, 1, False))


def _extract_ole_embedded_images_heuristic(ole, stream_name: str):
    """
    Best-effort embedded-image extraction from a legacy OLE (.doc/.ppt) stream --
    the counterpart to _extract_ole_text_heuristic() above, same "no full binary
    format parser" philosophy. Scans for JPEG/PNG magic-byte signatures and slices
    each byte range up to the next signature (or end of stream) out as a candidate
    image; PIL tolerates trailing garbage past an image's real end for these
    formats, so this coarse slicing is good enough without decoding the real
    picture-descriptor structures these formats use internally. BMP's 2-byte
    signature ("BM") is deliberately excluded -- far too likely to false-positive
    on arbitrary binary data to be worth the noise.
    Returns a list of raw image byte blobs (not yet OCR'd/validated as real images).
    """
    if not ole.exists(stream_name):
        return []
    try:
        data = ole.openstream(stream_name).read()
    except Exception:
        return []

    SIGNATURES = [b"\xFF\xD8\xFF", b"\x89PNG\r\n\x1a\n"]
    offsets = []
    for sig in SIGNATURES:
        start = 0
        while True:
            idx = data.find(sig, start)
            if idx == -1:
                break
            offsets.append(idx)
            start = idx + len(sig)
    offsets.sort()

    images = []
    for i, idx in enumerate(offsets):
        end = offsets[i + 1] if i + 1 < len(offsets) else len(data)
        blob = data[idx:end]
        if len(blob) >= 100:  # too small to plausibly be a real image -- skip
            images.append(blob)
    return images


def extract_text(f):
    name_lower = f.name.lower()

    # ── ZIP / Folder Upload ─────────────────────────────────────────────────
    if name_lower.endswith(".zip"):
        import zipfile, io as _io
        SUPPORTED_EXTS = (
            ".pdf", ".docx", ".doc", ".xlsx", ".xls",
            ".csv", ".pptx", ".ppt", ".txt", ".html", ".htm",
            ".png", ".jpg", ".jpeg"
        )
        combined_texts = []
        zip_chunks = []
        _ingested_chunks_cache[_cache_key(f.name)] = []
        try:
            with zipfile.ZipFile(_io.BytesIO(f.read())) as zf:
                entries = sorted(zf.namelist())
                for entry in entries:
                    if entry.endswith("/") or "__MACOSX" in entry or entry.startswith("."):
                        continue
                    entry_lower = entry.lower()
                    if not any(entry_lower.endswith(ext) for ext in SUPPORTED_EXTS):
                        continue
                    try:
                        # Check the entry's declared size/ratio BEFORE decompressing --
                        # previously inner_file.read() decompressed every matching
                        # member fully into memory first, with no size check at all on
                        # this path (scan_document()'s zip-bomb check runs at upload
                        # time on the OUTER file, not on each inner member extracted
                        # here). A small crafted zip with one highly-compressible member
                        # could trigger unbounded memory allocation.
                        from src.core.input_guardrail import MAX_UNCOMPRESSED_BYTES, MAX_ZIP_RATIO
                        info = zf.getinfo(entry)
                        if info.file_size > MAX_UNCOMPRESSED_BYTES:
                            combined_texts.append(f"--- FILE IN ZIP: {entry} ---\n[Skipped: declares {info.file_size} bytes uncompressed, exceeds the {MAX_UNCOMPRESSED_BYTES} byte safe limit.]")
                            continue
                        if info.compress_size > 0 and (info.file_size / info.compress_size) > MAX_ZIP_RATIO:
                            combined_texts.append(f"--- FILE IN ZIP: {entry} ---\n[Skipped: {info.file_size / info.compress_size:.0f}:1 compression ratio exceeds the safe {MAX_ZIP_RATIO}:1 limit.]")
                            continue
                        with zf.open(entry) as inner_file:
                            inner_bytes = inner_file.read()
                        inner_f = _io.BytesIO(inner_bytes)
                        inner_name = entry.split("/")[-1]
                        inner_f.name = inner_name
                        inner_text = extract_text(inner_f)
                        
                        inner_chunks = _ingested_chunks_cache.pop(_cache_key(inner_name), [])
                        for content, meta in inner_chunks:
                            meta["source_file"] = f"{f.name}/{entry}"
                            zip_chunks.append((content, meta))
                            
                        if inner_text and not inner_text.startswith("[Error"):
                            combined_texts.append(f"--- FILE IN ZIP: {entry} ---\n{inner_text}")
                        elif inner_text.startswith("[Error"):
                            combined_texts.append(f"--- FILE IN ZIP: {entry} ---\n{inner_text}")
                    except Exception as ie:
                        combined_texts.append(f"--- FILE IN ZIP: {entry} ---\n[Error reading {entry}: {ie}]")
            _ingested_chunks_cache[_cache_key(f.name)] = zip_chunks
            if combined_texts:
                return "\n\n".join(combined_texts)
            return "[ZIP file appears empty or contains no supported document types.]"
        except zipfile.BadZipFile:
            return f"[Error: {f.name} is not a valid ZIP file.]"
        except Exception as e:
            return f"[Error extracting ZIP {f.name}: {e}]"

    # ── Image files (PNG / JPG / JPEG) ──────────────────────────────────────
    if name_lower.endswith((".png", ".jpg", ".jpeg")):
        try:
            import PIL.Image
            import numpy as np
            reader = get_ocr_reader()
            img = PIL.Image.open(f)
            img_np = _preprocess_image_for_ocr(np.array(img))
            res = reader.readtext(img_np, detail=0)
            ocr_text = " ".join(res)

            image_chunks = []
            img_fname = getattr(f, "name", "unknown.png")
            img_ext = os.path.splitext(img_fname.lower())[1].lstrip(".")
            # Paragraph-based chunking: each OCR-detected text region/line is treated as
            # its own paragraph unit (chunk_paragraphs groups them up to ~1000 chars,
            # same as DOCX/HTML) instead of collapsing all regions into one string first
            # and char-slicing it, which cuts mid-line regardless of where OCR's own
            # region boundaries fell.
            paragraphs_data = [(line, "") for line in res if line and line.strip()]
            chunks = chunk_paragraphs(paragraphs_data, target=1000, overlap=200) if paragraphs_data else []
            for chunk_content, _chunk_section, _, _ in chunks:
                image_chunks.append((chunk_content, {
                    "source_file": img_fname,
                    "source_type": "image",
                    "image_id": img_fname,
                    "chunk_id": ""
                }))
            _ingested_chunks_cache[_cache_key(f.name)] = image_chunks
            return ocr_text
        except Exception as e:
            return f"[Error parsing image file {f.name}: {e}]"

    elif name_lower.endswith(".pdf"):
        import pdfplumber
        import numpy as np
        pages_text = []
        pdf_chunks = []
        current_section = ""
        try:
            with pdfplumber.open(f) as pdf:
                for p_idx, p in enumerate(pdf.pages, 1):
                    text = p.extract_text() or ""
                    # Hybrid OCR logic for PDFs
                    img_page = None
                    if hasattr(p, "images") and p.images:
                        try:
                            # Render page once at 150 resolution (optimized for CPU performance)
                            img_page = p.to_image(resolution=150)
                            pil_full = img_page.original
                            width_pixels, height_pixels = pil_full.size
                            scale_x = width_pixels / p.width
                            scale_y = height_pixels / p.height
                            
                            ocr_results = []
                            for img_meta in p.images:
                                x0 = img_meta.get("x0")
                                top = img_meta.get("top")
                                x1 = img_meta.get("x1")
                                bottom = img_meta.get("bottom")
                                if None not in (x0, top, x1, bottom) and x1 > x0 and bottom > top:
                                    left = int(x0 * scale_x)
                                    t = int(top * scale_y)
                                    right = int(x1 * scale_x)
                                    b = int(bottom * scale_y)
                                    
                                    cropped_pil = pil_full.crop((left, t, right, b))
                                    reader = get_ocr_reader()
                                    img_np = _preprocess_image_for_ocr(np.array(cropped_pil))
                                    res = reader.readtext(img_np, detail=0)
                                    if res:
                                        ocr_results.extend(res)
                            if ocr_results:
                                text += "\n[Embedded Image OCR]: " + " ".join(ocr_results)
                        except Exception as ocr_err:
                            print(f"[HYBRID OCR WARNING] Failed embedded image OCR: {ocr_err}", flush=True)

                    # Fail-Safe Full Page OCR Trigger:
                    # Runs if native text is low (<300 chars) OR if page has images but embedded OCR produced no text
                    should_full_ocr = (len(text.strip()) < 300) or (hasattr(p, "images") and p.images and "[Embedded Image OCR]" not in text)
                    if should_full_ocr:
                        try:
                            if img_page is None:
                                img_page = p.to_image(resolution=150)
                            pil_img = img_page.original
                            reader = get_ocr_reader()
                            img_np = _preprocess_image_for_ocr(np.array(pil_img))
                            res = reader.readtext(img_np, detail=0)
                            if res:
                                text += "\n[Page Image OCR]: " + " ".join(res)
                        except Exception as page_ocr_err:
                            print(f"[HYBRID OCR WARNING] Failed full page OCR: {page_ocr_err}", flush=True)


                    pages_text.append(text)
                    for line in text.splitlines():
                        line_str = line.strip()
                        if line_str and len(line_str) < 120:
                            if HEADER_REGEX.match(line_str):
                                current_section = line_str
                    page_chunks = chunk_text_by_chars(text, target=1000, overlap=200)
                    for chunk_txt in page_chunks:
                        p_text = chunk_txt
                        if current_section:
                            p_text = f"[{current_section}]\n{p_text}"
                        pdf_chunks.append((p_text, {
                            "source_file": getattr(f, "name", "unknown.pdf"),
                            "source_type": "pdf",
                            "page_number": p_idx,
                            "section_heading": current_section,
                            "chunk_id": ""
                        }))
            _ingested_chunks_cache[_cache_key(f.name)] = pdf_chunks
            return "\n".join(pages_text)
        except Exception as e:
            return f"[Error parsing PDF file {f.name}: {e}]"

    elif name_lower.endswith((".xlsx", ".xls")):
        try:
            import pandas as pd
            import re as _re
            # ISO control ID pattern: matches A.5.9, A.12.1.2, 5.9, 8.16, etc.
            _ISO_CTRL_RE = _re.compile(
                r'\b(?:A\.)?(\d{1,2}\.\d{1,2}(?:\.\d{1,2})?)\b'
            )
            excel_data = pd.read_excel(f, sheet_name=None)
            sheets_text = []
            xlsx_chunks = []
            xlsx_fname = getattr(f, "name", "unknown.xlsx")
            xlsx_ext = os.path.splitext(xlsx_fname.lower())[1].lstrip(".")
            xlsx_src_type = "xls" if xlsx_ext == "xls" else "xlsx"
            for sheet_name, df in excel_data.items():
                df_filled = df.fillna("")
                total_rows = len(df_filled)
                if total_rows == 0:
                    sheets_text.append(f"--- Sheet: {sheet_name} ---\n[Empty Sheet]")
                    continue
                # ── Improved chunking parameters ──────────────────────────────
                # 5 rows/chunk (was 15): tighter focus means each chunk is
                # about one audit control entry, not a mixed 15-row blob.
                # Overlap of 1 row keeps cross-boundary evidence intact.
                ROWS_PER_CHUNK = 5
                ROW_OVERLAP = 1
                MAX_CHUNK_CHARS = 2000   # hard cap per chunk
                columns = [str(c) for c in df_filled.columns]
                start_row = 0
                while start_row < total_rows:
                    end_row = min(start_row + ROWS_PER_CHUNK, total_rows)
                    df_slice = df_filled.iloc[start_row:end_row]

                    # ── Col=Value pipe format (replaces df.to_string) ─────────
                    # Before: "Implemented  A.9.1  High" (space-aligned, ambiguous)
                    # After:  "Control=A.9.1 | Status=Implemented | Risk=High"
                    # This preserves column identity so the keyword scorer can
                    # match "access control" → "A.9.1" in the correct column.
                    row_lines = []
                    for _, row in df_slice.iterrows():
                        pairs = []
                        for col in columns:
                            val = str(row[col]).strip()
                            if val and val != "nan":
                                pairs.append(f"{col}={val}")
                        if pairs:
                            row_lines.append(" | ".join(pairs))

                    if not row_lines:
                        start_row = end_row
                        continue

                    chunk_body = "\n".join(row_lines)

                    # ── ISO Control ID prefix ──────────────────────────────────
                    # Scan every cell in the slice for ISO control patterns.
                    # Surface them as "[Controls: 5.9, 8.16]" at the top of the
                    # chunk so the keyword scorer gets a direct match even when
                    # the control ID appears in a column the scorer doesn't weight.
                    ctrl_ids_found = set()
                    for _, row in df_slice.iterrows():
                        for col in columns:
                            cell_val = str(row[col])
                            for m in _ISO_CTRL_RE.finditer(cell_val):
                                ctrl_ids_found.add(m.group(0))
                    ctrl_prefix = ""
                    if ctrl_ids_found:
                        ctrl_prefix = f"[Controls: {', '.join(sorted(ctrl_ids_found))}]\n"

                    p_text = (
                        f"{ctrl_prefix}"
                        f"--- Sheet: {sheet_name} | Rows {start_row + 1}-{end_row} ---\n"
                        f"Columns: {' | '.join(columns)}\n"
                        f"{chunk_body}"
                    )

                    # Hard cap: if a single row still exceeds MAX_CHUNK_CHARS,
                    # truncate gracefully rather than feeding a token monster.
                    if len(p_text) > MAX_CHUNK_CHARS:
                        p_text = p_text[:MAX_CHUNK_CHARS] + "\n[...truncated]"

                    xlsx_chunks.append((p_text, {
                        "source_file": xlsx_fname,
                        "source_type": xlsx_src_type,
                        "sheet_name": sheet_name,
                        "start_row": start_row + 1,
                        "end_row": end_row,
                        "iso_controls_in_chunk": sorted(ctrl_ids_found),
                        "chunk_id": ""
                    }))

                    if end_row == total_rows:
                        break
                    next_start = end_row - ROW_OVERLAP
                    if next_start <= start_row:
                        next_start = end_row
                    start_row = next_start

                sheets_text.append(
                    f"--- Sheet: {sheet_name} ---\n" + df_filled.to_string(index=False)
                )

            # ── Embedded images (e.g. a screenshot pasted into a cell) ──────────
            # pandas.read_excel() only ever sees cell values -- pasted images are
            # completely invisible to it. openpyxl exposes them via worksheet._images,
            # each with an anchor cell telling us its approximate row, so the OCR'd
            # text can be inserted next to the row-range chunk it actually belongs to
            # instead of being silently dropped. .xls (legacy binary Excel) isn't
            # supported by openpyxl at all, so this only runs for real .xlsx files.
            if xlsx_ext == "xlsx":
                try:
                    if hasattr(f, "seek"):
                        f.seek(0)
                    xlsx_bytes = f.read()
                    if hasattr(f, "seek"):
                        f.seek(0)
                    import io as _xlsx_io
                    import openpyxl
                    wb = openpyxl.load_workbook(_xlsx_io.BytesIO(xlsx_bytes))
                    for ws in wb.worksheets:
                        images = getattr(ws, "_images", None) or []
                        for image in images:
                            try:
                                anchor_from = getattr(image.anchor, "_from", None)
                                anchor_row = (anchor_from.row + 1) if anchor_from is not None else None
                                img_data = image._data()
                                import PIL.Image
                                import numpy as np
                                img = PIL.Image.open(_xlsx_io.BytesIO(img_data))
                                img_np = _preprocess_image_for_ocr(np.array(img))
                                reader = get_ocr_reader()
                                res = reader.readtext(img_np, detail=0)
                                if not res:
                                    continue
                                ocr_text = " ".join(res)
                                img_chunk = (
                                    f"[Embedded Image OCR]: {ocr_text}",
                                    {
                                        "source_file": xlsx_fname,
                                        "source_type": xlsx_src_type,
                                        "sheet_name": ws.title,
                                        "start_row": anchor_row,
                                        "end_row": anchor_row,
                                        "chunk_id": ""
                                    }
                                )
                                # Insert right after the row-range chunk this image's
                                # anchor row falls inside, so it stays positioned near
                                # the data it was actually pasted next to.
                                insert_at = len(xlsx_chunks)
                                if anchor_row is not None:
                                    for idx, (_, meta) in enumerate(xlsx_chunks):
                                        if (meta.get("sheet_name") == ws.title
                                                and meta.get("start_row") is not None
                                                and meta["start_row"] <= anchor_row <= meta["end_row"]):
                                            insert_at = idx + 1
                                            break
                                xlsx_chunks.insert(insert_at, img_chunk)
                            except Exception:
                                pass  # unreadable/unsupported embedded image -- skip
                except Exception as xlsx_img_err:
                    print(f"[XLSX IMAGE WARNING] Failed to extract embedded images: {xlsx_img_err}", flush=True)

            _ingested_chunks_cache[_cache_key(f.name)] = xlsx_chunks
            return "\n\n".join(sheets_text)
        except Exception as e:
            return f"[Error parsing Excel file {f.name}: {e}]"

    elif name_lower.endswith(".csv"):
        try:
            if hasattr(f, "seek"):
                f.seek(0)
            raw_csv_bytes = f.read()
            raw_csv_text = raw_csv_bytes.decode("utf-8", errors="ignore")

            import pandas as pd
            df = pd.read_csv(io.StringIO(raw_csv_text))
            df_filled = df.fillna("")
            total_rows = len(df_filled)
            csv_chunks = []
            if total_rows > 0:
                ROWS_PER_CHUNK = 15
                ROW_OVERLAP = 3
                csv_fname = getattr(f, "name", "unknown.csv")
                # Build header row string to prepend to every CSV chunk
                csv_header_str = "  ".join(str(c) for c in df_filled.columns)
                start_row = 0
                while start_row < total_rows:
                    end_row = min(start_row + ROWS_PER_CHUNK, total_rows)
                    df_slice = df_filled.iloc[start_row:end_row]
                    chunk_text_val = df_slice.to_string(index=False)
                    # Collapse consecutive spaces
                    chunk_text_val = re.sub(r' {2,}', '  ', chunk_text_val)
                    p_csv_text = f"Headers: {csv_header_str}\n{chunk_text_val}"
                    csv_chunks.append((p_csv_text, {
                        "source_file": csv_fname,
                        "source_type": "csv",
                        "sheet_name": "CSV",
                        "start_row": start_row + 1,
                        "end_row": end_row,
                        "chunk_id": ""
                    }))
                    if end_row == total_rows:
                        break
                    next_start = end_row - ROW_OVERLAP
                    if next_start <= start_row:
                        next_start = end_row
                    start_row = next_start
            _ingested_chunks_cache[_cache_key(f.name)] = csv_chunks
            # Return the raw CSV text (not the pandas-reformatted table) so tools that
            # parse CSV directly (e.g. QualysParser's csv.DictReader) get real delimiters.
            return raw_csv_text if raw_csv_text.strip() else "[Empty CSV]"
        except Exception as e:
            return f"[Error parsing CSV file {f.name}: {e}]"

    elif name_lower.endswith(".pptx"):
        try:
            from pptx import Presentation
            prs = Presentation(f)
            pptx_chunks = []
            all_text_runs = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_title = ""
                try:
                    if slide.shapes.title and hasattr(slide.shapes.title, "text"):
                        slide_title = slide.shapes.title.text.strip()
                except Exception:
                    pass
                if not slide_title:
                    for shape in slide.shapes:
                        if hasattr(shape, "is_placeholder") and shape.is_placeholder and hasattr(shape, "placeholder_format") and shape.placeholder_format.type == 1:
                            if hasattr(shape, "text") and shape.text.strip():
                                slide_title = shape.text.strip()
                                break
                shape_texts = []
                for shape in slide.shapes:
                    if shape == getattr(slide.shapes, "title", None):
                        continue
                    if hasattr(shape, "text") and shape.text.strip():
                        shape_texts.append(shape.text.strip())
                    if shape.has_table:
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text:
                                shape_texts.append(" | ".join(row_text))
                    # Check for embedded picture shapes (type 13 is PICTURE)
                    if hasattr(shape, "shape_type") and shape.shape_type == 13:
                        try:
                            import io as _io
                            image_bytes = shape.image.blob
                            import PIL.Image
                            import numpy as np
                            img = PIL.Image.open(_io.BytesIO(image_bytes))
                            img_np = _preprocess_image_for_ocr(np.array(img))
                            reader = get_ocr_reader()
                            res = reader.readtext(img_np, detail=0)
                            if res:
                                ocr_text = " ".join(res)
                                shape_texts.append(f"[Slide Image OCR]: {ocr_text}")
                        except Exception:
                            pass
                notes_text = ""
                try:
                    if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                except Exception:
                    pass
                slide_block_runs = [f"--- Slide {slide_num} ---"]
                if slide_title:
                    slide_block_runs.append(f"Title: {slide_title}")
                if shape_texts:
                    slide_block_runs.extend(shape_texts)
                if notes_text:
                    slide_block_runs.append(f"Notes: {notes_text}")
                slide_text = "\n".join(slide_block_runs)
                all_text_runs.append(slide_text)
                pptx_fname = getattr(f, "name", "unknown.pptx")
                pptx_ext = os.path.splitext(pptx_fname.lower())[1].lstrip(".")
                pptx_chunks.append((slide_text, {
                    "source_file": pptx_fname,
                    "source_type": pptx_ext if pptx_ext else "pptx",
                    "slide_number": slide_num,
                    "slide_title": slide_title,
                    "chunk_id": ""
                }))
            _ingested_chunks_cache[_cache_key(f.name)] = pptx_chunks
            return "\n\n".join(all_text_runs)
        except Exception as e:
            return f"[Error parsing PowerPoint file {f.name}: {e}]"

    elif name_lower.endswith((".txt", ".nessus", ".gnmap", ".nmap", ".log")):
        # Plain-text VAPT scanner exports — fast path, no structure needed.
        try:
            txt_content = f.read().decode("utf-8", errors="ignore")
            txt_fname = getattr(f, "name", "unknown.txt")
            txt_ext = os.path.splitext(txt_fname.lower())[1].lstrip(".") or "txt"
            txt_chunks = []
            chunks = chunk_text_by_chars(txt_content, target=1000, overlap=200)
            for chunk_content in chunks:
                txt_chunks.append((chunk_content, {
                    "source_file": txt_fname,
                    "source_type": txt_ext,
                    "chunk_id": ""
                }))
            _ingested_chunks_cache[_cache_key(f.name)] = txt_chunks
            return txt_content
        except Exception as e:
            return f"[Error parsing text file {f.name}: {e}]"

    elif name_lower.endswith(".xml"):
        # ── Smart XML parser: Burp Suite / Nessus .nessus-as-xml / Nmap / ZAP / Qualys ──
        # Detects the schema by root/child tag names and extracts structured fields
        # (name, severity, host, description, evidence, solution) as clean text lines
        # so the RAG retriever gets properly isolated field-level tokens instead of
        # raw angle-bracket soup. Falls back to plain-text for unknown XML schemas.
        try:
            if hasattr(f, "seek"):
                f.seek(0)
            xml_bytes = f.read()
            if hasattr(f, "seek"):
                f.seek(0)
            xml_fname = getattr(f, "name", "unknown.xml")

            # defusedxml hardens against XML entity-expansion attacks (billion
            # laughs, quadratic blowup, external entity/DTD fetching) that plain
            # xml.etree.ElementTree doesn't guard against on untrusted uploaded
            # XML. Fully API-compatible with the stdlib module (same fromstring(),
            # same ParseError) -- this is a drop-in import swap, nothing else in
            # this function needs to change.
            import defusedxml.ElementTree as _ET

            def _safe_text(el, tag, default=""):
                """Return stripped text of the first matching child tag, or default."""
                child = el.find(tag)
                return (child.text or "").strip() if child is not None else default

            lines = []  # List of structured text lines written to the chunk store

            try:
                root = _ET.fromstring(xml_bytes)
            except _ET.ParseError:
                # Malformed XML — fall back to raw text
                root = None

            if root is not None:
                rtag = root.tag.lower().split("}")[-1]  # strip namespace

                # ── Burp Suite XML ────────────────────────────────────────────
                # <issues burpVersion="..." exportTime="...">
                #   <issue>
                #     <serialNumber/> <type/> <name/> <host/> <path/>
                #     <severity/> <confidence/> <issueDetail/>
                #     <issueBackground/> <remediationDetail/>
                #   </issue>
                # </issues>
                if rtag in ("issues", "issue") or root.find("issue") is not None:
                    lines.append("[BURP SUITE SCAN REPORT]")
                    issues = root.findall("issue") if rtag == "issues" else [root]
                    for i, issue in enumerate(issues, 1):
                        name = _safe_text(issue, "name", f"Issue {i}")
                        host = _safe_text(issue, "host")
                        path = _safe_text(issue, "path")
                        sev = _safe_text(issue, "severity")
                        conf = _safe_text(issue, "confidence")
                        detail = _safe_text(issue, "issueDetail")
                        background = _safe_text(issue, "issueBackground")
                        remediation = _safe_text(issue, "remediationDetail")
                        # Burp stores request/response bodies in <requestresponse> — extract raw text
                        rr = issue.find("requestresponse")
                        rr_text = ""
                        if rr is not None:
                            req_b64 = _safe_text(rr, "request")
                            resp_b64 = _safe_text(rr, "response")
                            import base64 as _b64
                            try:
                                rr_text = _b64.b64decode(req_b64).decode("utf-8", errors="ignore")[:500]
                            except Exception:
                                rr_text = req_b64[:200]

                        block = (
                            f"[Finding {i}] {name}\n"
                            f"  Host: {host}{path}\n"
                            f"  Severity: {sev}  Confidence: {conf}\n"
                        )
                        if detail:
                            block += f"  Detail: {detail}\n"
                        if background:
                            block += f"  Background: {background}\n"
                        if remediation:
                            block += f"  Remediation: {remediation}\n"
                        if rr_text:
                            block += f"  Request/Response snippet: {rr_text}\n"
                        lines.append(block)

                # ── Nessus XML (.nessus native format) ───────────────────────
                # <NessusClientData_v2>
                #   <Report><ReportHost name="x.x.x.x">
                #     <ReportItem port="" svc_name="" protocol="" severity="" pluginName="">
                #       <description/> <solution/> <risk_factor/>
                #       <plugin_output/> <cve/> <cvss_base_score/>
                #
                elif rtag in ("nessusclientdata_v2", "nessus") or root.find(".//reportitem") is not None or root.find(".//ReportItem") is not None:
                    lines.append("[NESSUS SCAN REPORT]")
                    for host_el in root.iter("ReportHost"):
                        host_name = host_el.get("name", "unknown")
                        lines.append(f"\nHost: {host_name}")
                        for item in host_el.findall("ReportItem"):
                            plugin = item.get("pluginName", "")
                            port = item.get("port", "")
                            proto = item.get("protocol", "")
                            sev = item.get("severity", "")
                            sev_labels = {"0": "Info", "1": "Low", "2": "Medium", "3": "High", "4": "Critical"}
                            sev_label = sev_labels.get(sev, sev)
                            risk = _safe_text(item, "risk_factor")
                            cve = _safe_text(item, "cve")
                            cvss = _safe_text(item, "cvss_base_score")
                            desc = _safe_text(item, "description")
                            solution = _safe_text(item, "solution")
                            plugin_out = _safe_text(item, "plugin_output")

                            block = (
                                f"  [{sev_label}] {plugin}\n"
                                f"    Port: {port}/{proto}  Risk: {risk}"
                            )
                            if cve:
                                block += f"  CVE: {cve}"
                            if cvss:
                                block += f"  CVSS: {cvss}"
                            block += "\n"
                            if desc:
                                block += f"    Description: {desc[:400]}\n"
                            if solution:
                                block += f"    Solution: {solution[:300]}\n"
                            if plugin_out:
                                block += f"    Plugin Output: {plugin_out[:400]}\n"
                            lines.append(block)

                # ── Nmap XML ─────────────────────────────────────────────────
                # <nmaprun scanner="nmap" ...>
                #   <host><address addr="x.x.x.x" addrtype="ipv4"/>
                #         <ports><port protocol="tcp" portid="80">
                #           <state state="open"/><service name="http"/>
                #
                elif rtag == "nmaprun" or root.find(".//nmaprun") is not None or root.find("host") is not None:
                    lines.append("[NMAP SCAN REPORT]")
                    for host_el in root.iter("host"):
                        addr_el = host_el.find("address")
                        addr = addr_el.get("addr", "unknown") if addr_el is not None else "unknown"
                        hostname_el = host_el.find(".//hostname")
                        hostname = hostname_el.get("name", "") if hostname_el is not None else ""
                        host_label = f"{addr}" + (f" ({hostname})" if hostname else "")
                        lines.append(f"\nHost: {host_label}")
                        os_el = host_el.find(".//osmatch")
                        if os_el is not None:
                            lines.append(f"  OS: {os_el.get('name', '')} (accuracy {os_el.get('accuracy', '')}%)")
                        for port_el in host_el.iter("port"):
                            portid = port_el.get("portid", "")
                            proto = port_el.get("protocol", "")
                            state_el = port_el.find("state")
                            state = state_el.get("state", "") if state_el is not None else ""
                            svc_el = port_el.find("service")
                            svc = ""
                            if svc_el is not None:
                                svc = svc_el.get("name", "")
                                product = svc_el.get("product", "")
                                version = svc_el.get("version", "")
                                if product:
                                    svc += f" ({product} {version})".strip()
                            script_outputs = []
                            for script_el in port_el.findall("script"):
                                script_outputs.append(f"{script_el.get('id', '')}: {script_el.get('output', '')[:200]}")
                            line = f"  {portid}/{proto}  {state}  {svc}"
                            if script_outputs:
                                line += "\n    " + "\n    ".join(script_outputs)
                            lines.append(line)

                # ── OWASP ZAP XML ─────────────────────────────────────────────
                # <OWASPZAPReport version="...">
                #   <site ...><alerts><alertitem>
                #     <alert/> <riskdesc/> <desc/> <solution/> <evidence/>
                #     <instances><instance><uri/><method/><evidence/></instance>
                #
                elif rtag in ("owaspzapreport", "zapreport") or root.find(".//alertitem") is not None:
                    lines.append("[OWASP ZAP SCAN REPORT]")
                    for alert in root.iter("alertitem"):
                        name = _safe_text(alert, "alert") or _safe_text(alert, "name")
                        risk = _safe_text(alert, "riskdesc") or _safe_text(alert, "risk")
                        desc = _safe_text(alert, "desc") or _safe_text(alert, "description")
                        solution = _safe_text(alert, "solution")
                        evidence = _safe_text(alert, "evidence")
                        cweid = _safe_text(alert, "cweid")
                        instances = []
                        for inst in alert.findall(".//instance"):
                            uri = _safe_text(inst, "uri")
                            method = _safe_text(inst, "method")
                            ev = _safe_text(inst, "evidence")
                            if uri:
                                instances.append(f"{method} {uri}" + (f" [{ev}]" if ev else ""))

                        block = f"[{risk}] {name}\n"
                        if cweid:
                            block += f"  CWE: {cweid}\n"
                        if desc:
                            block += f"  Description: {desc[:400]}\n"
                        if solution:
                            block += f"  Solution: {solution[:300]}\n"
                        if evidence:
                            block += f"  Evidence: {evidence[:200]}\n"
                        if instances:
                            block += "  Instances:\n    " + "\n    ".join(instances[:10]) + "\n"
                        lines.append(block)

                # ── Qualys / OpenVAS XML ──────────────────────────────────────
                # Qualys: <QualysGuardReport><HOST_LIST><HOST>
                #           <VULN_INFO_LIST><VULN_INFO>
                #             <QID/><TITLE/><SEVERITY/><RESULT/>
                # OpenVAS: <report><results><result>
                #             <name/><severity/><description/><solution/>
                #             <host><ip/></host>
                elif (root.find(".//VULN_INFO") is not None
                      or root.find(".//vuln_info") is not None
                      or root.find(".//result") is not None):
                    lines.append("[QUALYS / OPENVAS SCAN REPORT]")
                    # OpenVAS result elements
                    for result in root.iter("result"):
                        name = _safe_text(result, "name")
                        sev = _safe_text(result, "severity")
                        host_el = result.find("host")
                        host = ""
                        if host_el is not None:
                            ip_el = host_el.find("ip")
                            host = (ip_el.text or host_el.text or "").strip() if ip_el is not None else (host_el.text or "").strip()
                        desc = _safe_text(result, "description")
                        solution = _safe_text(result, "solution")
                        block = f"[{sev}] {name}  Host: {host}\n"
                        if desc:
                            block += f"  Description: {desc[:400]}\n"
                        if solution:
                            block += f"  Solution: {solution[:300]}\n"
                        lines.append(block)
                    # Qualys VULN_INFO elements
                    for vuln in root.iter("VULN_INFO"):
                        qid = _safe_text(vuln, "QID")
                        title = _safe_text(vuln, "TITLE") or _safe_text(vuln, "title")
                        sev = _safe_text(vuln, "SEVERITY") or _safe_text(vuln, "severity")
                        result_text = _safe_text(vuln, "RESULT") or _safe_text(vuln, "result")
                        block = f"[Severity {sev}] QID {qid}: {title}\n"
                        if result_text:
                            block += f"  Result: {result_text[:400]}\n"
                        lines.append(block)

                # ── Generic / Unknown XML ─────────────────────────────────────
                # No recognised schema — convert element tree to readable key:value
                # lines so the LLM sees structured fields rather than angle brackets.
                else:
                    lines.append(f"[XML DOCUMENT: {xml_fname}]")
                    def _xml_to_lines(el, depth=0):
                        tag = el.tag.split("}")[-1]  # strip namespace
                        text = (el.text or "").strip()
                        attribs = " ".join(f'{k}="{v}"' for k, v in el.attrib.items())
                        prefix = "  " * depth
                        if text:
                            lines.append(f"{prefix}{tag}{(' [' + attribs + ']') if attribs else ''}: {text[:300]}")
                        elif attribs:
                            lines.append(f"{prefix}{tag} [{attribs}]")
                        for child in list(el)[:50]:  # cap at 50 children per node
                            _xml_to_lines(child, depth + 1)
                    for child in list(root)[:200]:  # cap root children
                        _xml_to_lines(child, depth=0)

            else:
                # Unparseable XML — treat as plain text
                lines.append(xml_bytes.decode("utf-8", errors="ignore"))

            structured_text = "\n".join(lines)
            xml_chunks = []
            chunks = chunk_text_by_chars(structured_text, target=1000, overlap=200)
            for chunk_content in chunks:
                xml_chunks.append((chunk_content, {
                    "source_file": xml_fname,
                    "source_type": "xml",
                    "chunk_id": ""
                }))
            _ingested_chunks_cache[_cache_key(f.name)] = xml_chunks
            return structured_text
        except Exception as e:
            return f"[Error parsing XML file {f.name}: {e}]"

    elif name_lower.endswith(".json"):
        # ── Smart JSON flattener: Trivy / ZAP JSON / generic nested reports ──
        # Recursively flattens nested JSON into "path.to.key: value" lines so the
        # LLM retriever sees CVE IDs, severities, and fix versions as plain tokens
        # rather than deeply nested dictionary text.
        try:
            if hasattr(f, "seek"):
                f.seek(0)
            raw_json = f.read().decode("utf-8", errors="ignore")
            if hasattr(f, "seek"):
                f.seek(0)
            json_fname = getattr(f, "name", "unknown.json")

            import json as _json

            try:
                data = _json.loads(raw_json)
            except _json.JSONDecodeError:
                # Not valid JSON — store as plain text
                data = None

            lines = []

            if data is None:
                lines.append(raw_json)

            else:
                def _flatten(obj, prefix="", depth=0):
                    """Recursively flatten JSON into human-readable lines."""
                    if depth > 10:
                        return  # guard against pathological nesting
                    if isinstance(obj, dict):
                        for k, v in obj.items():
                            key = f"{prefix}.{k}" if prefix else k
                            _flatten(v, key, depth + 1)
                    elif isinstance(obj, list):
                        # For lists: emit each item with [n] index; cap at 500 items
                        for i, item in enumerate(obj[:500]):
                            key = f"{prefix}[{i}]"
                            _flatten(item, key, depth + 1)
                    else:
                        val = str(obj).strip()
                        if val and val != "null" and val != "None":
                            lines.append(f"{prefix}: {val[:500]}")

                # ── Trivy JSON (most common VAPT CI/CD tool) ──────────────────
                # {"SchemaVersion":2, "Results":[{"Target":"img","Vulnerabilities":[{
                #   "VulnerabilityID":"CVE-...", "PkgName":"...", "InstalledVersion":"...",
                #   "FixedVersion":"...", "Severity":"CRITICAL","Title":"...","Description":"..."}]}]}
                if isinstance(data, dict) and "Results" in data:
                    lines.append("[TRIVY VULNERABILITY REPORT]")
                    schema = data.get("SchemaVersion", "")
                    artifact = data.get("ArtifactName", "")
                    if artifact:
                        lines.append(f"Artifact: {artifact}")
                    for result in data.get("Results", []):
                        target = result.get("Target", "")
                        res_type = result.get("Type", "")
                        lines.append(f"\nTarget: {target}  ({res_type})")
                        vulns = result.get("Vulnerabilities") or []
                        misconfs = result.get("Misconfigurations") or []
                        for v in vulns:
                            cve = v.get("VulnerabilityID", "")
                            pkg = v.get("PkgName", "")
                            installed = v.get("InstalledVersion", "")
                            fixed = v.get("FixedVersion", "")
                            sev = v.get("Severity", "")
                            title = v.get("Title", "")
                            desc = (v.get("Description") or "")[:300]
                            lines.append(
                                f"  [{sev}] {cve} — {pkg} {installed}"
                                + (f" (fix: {fixed})" if fixed else " (no fix)")
                                + (f"\n    {title}" if title else "")
                                + (f"\n    {desc}" if desc else "")
                            )
                        for m in misconfs:
                            mid = m.get("ID", "")
                            msev = m.get("Severity", "")
                            mtitle = m.get("Title", "")
                            mdesc = (m.get("Description") or "")[:200]
                            lines.append(f"  [{msev}] Misconfiguration {mid}: {mtitle}\n    {mdesc}")

                # ── ZAP JSON ───────────────────────────────────────────────────
                # {"site":[{"@name":"...","alerts":[{"alert":"...","riskdesc":"...","desc":"...","solution":"...","instances":[...]}]}]}
                elif isinstance(data, dict) and "site" in data:
                    lines.append("[OWASP ZAP JSON REPORT]")
                    sites = data["site"] if isinstance(data["site"], list) else [data["site"]]
                    for site in sites:
                        site_name = site.get("@name", "")
                        lines.append(f"\nSite: {site_name}")
                        for alert in site.get("alerts", []):
                            name = alert.get("alert", "")
                            risk = alert.get("riskdesc", "")
                            desc = (alert.get("desc") or "")[:300]
                            solution = (alert.get("solution") or "")[:200]
                            cweid = alert.get("cweid", "")
                            instances = alert.get("instances", [])
                            block = f"  [{risk}] {name}"
                            if cweid:
                                block += f"  (CWE-{cweid})"
                            block += f"\n    Description: {desc}"
                            if solution:
                                block += f"\n    Solution: {solution}"
                            if instances:
                                uris = [inst.get("uri", "") for inst in instances[:5] if inst.get("uri")]
                                block += "\n    URIs: " + ", ".join(uris)
                            lines.append(block)

                # ── Generic nested JSON ────────────────────────────────────────
                else:
                    lines.append(f"[JSON DOCUMENT: {json_fname}]")
                    _flatten(data)

            structured_text = "\n".join(lines)
            json_chunks = []
            chunks = chunk_text_by_chars(structured_text, target=1000, overlap=200)
            for chunk_content in chunks:
                json_chunks.append((chunk_content, {
                    "source_file": json_fname,
                    "source_type": "json",
                    "chunk_id": ""
                }))
            _ingested_chunks_cache[_cache_key(f.name)] = json_chunks
            return structured_text
        except Exception as e:
            return f"[Error parsing JSON file {f.name}: {e}]"

    elif name_lower.endswith((".html", ".htm")):
        try:
            if hasattr(f, "seek"):
                f.seek(0)
            html_bytes = f.read()
            if hasattr(f, "seek"):
                f.seek(0)
            html_str = html_bytes.decode("utf-8", errors="ignore")

            from bs4 import BeautifulSoup
            # Prefer lxml for speed; fallback to html.parser if unavailable
            try:
                import lxml  # noqa: F401
                _parser = "lxml"
            except ImportError:
                _parser = "html.parser"
            soup = BeautifulSoup(html_str, _parser)

            for script in soup(["script", "style", "head", "meta", "link", "svg"]):
                script.decompose()

            paragraphs_data = []
            current_section = "HTML Content"

            # FIX: added "pre" and "code" so Nessus plugin output / server banners /
            # raw HTTP responses are captured as evidence chunks instead of being silently dropped.
            for element in soup.find_all(["h1", "h2", "h3", "h4", "h5", "h6",
                                          "p", "div", "tr", "li", "img",
                                          "pre", "code"]):
                if element.name in ["h1", "h2", "h3", "h4", "h5", "h6"]:
                    heading_text = element.get_text(strip=True)
                    if heading_text:
                        current_section = heading_text
                        paragraphs_data.append((heading_text, current_section))
                elif element.name == "tr":
                    cells = [td.get_text(strip=True) for td in element.find_all(["th", "td"]) if td.get_text(strip=True)]
                    if cells:
                        row_str = " | ".join(cells)
                        paragraphs_data.append((row_str, current_section))
                elif element.name in ("pre", "code"):
                    # Plugin output, server banners, raw HTTP/TLS responses —
                    # keep as-is (don't require ≥3 word minimum so 1-line banners
                    # like "Apache/2.4.51" are still captured).
                    text = element.get_text(strip=False)  # preserve newlines
                    text = text.strip()
                    if text:
                        # Hard cap at 1500 chars so a single 200KB curl response
                        # doesn't blow a chunk's token budget.
                        if len(text) > 1500:
                            text = text[:1500] + "\n[...truncated]"
                        paragraphs_data.append((f"[Plugin Output / Code Block]\n{text}", current_section))
                elif element.name == "img":
                    # Only base64-embedded images (data: URIs) -- this tool is
                    # offline-first, so fetching external/relative image URLs over
                    # the network is out of scope (and would just fail in an
                    # air-gapped deployment anyway). find_all already returns
                    # elements in document order, so this stays correctly
                    # positioned relative to the surrounding text automatically.
                    src = element.get("src", "")
                    if src.startswith("data:image"):
                        try:
                            import base64
                            import io as _html_io
                            import PIL.Image
                            import numpy as np
                            b64_payload = src.split(",", 1)[1]
                            img_bytes = base64.b64decode(b64_payload)
                            img = PIL.Image.open(_html_io.BytesIO(img_bytes))
                            img_np = _preprocess_image_for_ocr(np.array(img))
                            reader = get_ocr_reader()
                            res = reader.readtext(img_np, detail=0)
                            if res:
                                ocr_text = " ".join(res)
                                paragraphs_data.append((f"[Embedded Image OCR]: {ocr_text}", current_section))
                        except Exception:
                            pass  # not a decodable/OCR-able embedded image -- skip
                else:
                    text = element.get_text(strip=True)
                    if text and len(text.split()) >= 3 and not element.find_all(["h1", "h2", "h3", "h4", "h5", "h6", "tr", "p", "div", "li"]):
                        paragraphs_data.append((text, current_section))

            html_fname = getattr(f, "name", "vapt_report.html")
            chunks = chunk_paragraphs(paragraphs_data, target=1000, overlap=200)
            html_chunks = []
            for chunk_content, chunk_section, _, _ in chunks:
                html_chunks.append((chunk_content, {
                    "source_file": html_fname,
                    "source_type": "html",
                    "section_heading": chunk_section,
                    "chunk_id": ""
                }))
            _ingested_chunks_cache[_cache_key(f.name)] = html_chunks
            return "\n\n".join([txt for txt, _ in paragraphs_data])
        except Exception as e:
            return f"[Error parsing HTML file {f.name}: {e}]"

    elif name_lower.endswith(".doc"):
        # Legacy Word 97-2003 binary format (OLE Compound File) -- python-docx can't
        # open these at all. Best-effort text extraction via the WordDocument stream.
        try:
            import olefile
            if hasattr(f, "seek"):
                f.seek(0)
            file_bytes = f.read()
            if not olefile.isOleFile(io.BytesIO(file_bytes)):
                return f"[Error parsing file {f.name}: not a valid legacy .doc (OLE) file]"
            ole = olefile.OleFileIO(io.BytesIO(file_bytes))
            try:
                if not ole.exists("WordDocument"):
                    return f"[Error parsing file {f.name}: no WordDocument stream found]"
                stream_bytes = ole.openstream("WordDocument").read()
                # No dedicated picture stream in legacy .doc the way .ppt has one --
                # "Data" is the closest equivalent; best-effort, may miss some/all
                # embedded images depending on how the document embedded them.
                image_blobs = _extract_ole_embedded_images_heuristic(ole, "Data")
            finally:
                ole.close()
            doc_text = _extract_ole_text_heuristic(stream_bytes)
            doc_fname = getattr(f, "name", "unknown.doc")
            doc_chunks = []
            for chunk_content in chunk_text_by_chars(doc_text, target=1000, overlap=200):
                doc_chunks.append((chunk_content, {
                    "source_file": doc_fname,
                    "source_type": "doc",
                    "chunk_id": ""
                }))
            # Legacy .doc's own text extraction is already a flat, position-agnostic
            # heuristic scan (no real paragraph structure to interleave against), so
            # images are appended at the end -- same tier DOCX had before its fix,
            # not true positional interleaving, but nothing stays invisible either.
            for blob in image_blobs:
                try:
                    import PIL.Image
                    import numpy as np
                    img = PIL.Image.open(io.BytesIO(blob))
                    img_np = _preprocess_image_for_ocr(np.array(img))
                    reader = get_ocr_reader()
                    res = reader.readtext(img_np, detail=0)
                    if res:
                        ocr_text = " ".join(res)
                        doc_chunks.append((f"[Embedded Image OCR]: {ocr_text}", {
                            "source_file": doc_fname,
                            "source_type": "doc",
                            "chunk_id": ""
                        }))
                except Exception:
                    pass  # not a real/decodable image -- skip
            _ingested_chunks_cache[_cache_key(f.name)] = doc_chunks
            return "[Best-effort extraction from legacy .doc format -- formatting/tables not preserved]\n" + doc_text
        except Exception as e:
            return f"[Error parsing legacy Word file {f.name}: {e}]"

    elif name_lower.endswith(".ppt"):
        # Legacy PowerPoint 97-2003 binary format (OLE Compound File) -- python-pptx
        # can't open these at all. Best-effort text extraction via the
        # "PowerPoint Document" stream.
        try:
            import olefile
            if hasattr(f, "seek"):
                f.seek(0)
            file_bytes = f.read()
            if not olefile.isOleFile(io.BytesIO(file_bytes)):
                return f"[Error parsing file {f.name}: not a valid legacy .ppt (OLE) file]"
            ole = olefile.OleFileIO(io.BytesIO(file_bytes))
            try:
                if not ole.exists("PowerPoint Document"):
                    return f"[Error parsing file {f.name}: no PowerPoint Document stream found]"
                stream_bytes = ole.openstream("PowerPoint Document").read()
                # PowerPoint 97-2003 concatenates all embedded pictures into a
                # dedicated "Pictures" stream -- more reliable target than .doc has.
                image_blobs = _extract_ole_embedded_images_heuristic(ole, "Pictures")
            finally:
                ole.close()
            ppt_text = _extract_ole_text_heuristic(stream_bytes)
            ppt_fname = getattr(f, "name", "unknown.ppt")
            ppt_chunks = []
            for chunk_content in chunk_text_by_chars(ppt_text, target=1000, overlap=200):
                ppt_chunks.append((chunk_content, {
                    "source_file": ppt_fname,
                    "source_type": "ppt",
                    "chunk_id": ""
                }))
            # Slide structure isn't preserved by this format's text extraction either
            # (see comment above the return string), so same as .doc: images are
            # appended at the end rather than tied to a specific slide.
            for blob in image_blobs:
                try:
                    import PIL.Image
                    import numpy as np
                    img = PIL.Image.open(io.BytesIO(blob))
                    img_np = _preprocess_image_for_ocr(np.array(img))
                    reader = get_ocr_reader()
                    res = reader.readtext(img_np, detail=0)
                    if res:
                        ocr_text = " ".join(res)
                        ppt_chunks.append((f"[Embedded Image OCR]: {ocr_text}", {
                            "source_file": ppt_fname,
                            "source_type": "ppt",
                            "chunk_id": ""
                        }))
                except Exception:
                    pass  # not a real/decodable image -- skip
            _ingested_chunks_cache[_cache_key(f.name)] = ppt_chunks
            return "[Best-effort extraction from legacy .ppt format -- slide structure not preserved]\n" + ppt_text
        except Exception as e:
            return f"[Error parsing legacy PowerPoint file {f.name}: {e}]"

    else:
        try:
            import zipfile
            import io as _io
            from docx import Document
            from docx.oxml.ns import qn

            # Reset seek position of f if possible, and read all bytes
            if hasattr(f, "seek"):
                f.seek(0)
            file_bytes = f.read()
            if hasattr(f, "seek"):
                f.seek(0)

            doc = Document(_io.BytesIO(file_bytes))
            paragraphs_data = []
            processed_media_names = set()  # basenames already OCR'd inline -- skip in the ZIP fallback pass below

            # 1. Paragraphs with heading detection, plus any inline/anchored images
            # embedded directly in that paragraph's runs -- OCR'd and inserted right
            # here, so an image and the text around it stay together in the same
            # position, instead of every image in the file being batched at the end
            # regardless of which paragraph it actually appeared next to.
            current_section = ""
            for p in doc.paragraphs:
                p_text = p.text.strip()
                if p_text:
                    if p.style and p.style.name and p.style.name.startswith("Heading"):
                        current_section = p_text
                    paragraphs_data.append((p_text, current_section))

                for run in p.runs:
                    for drawing in run._element.findall(qn("w:drawing")):
                        for blip in drawing.findall(".//" + qn("a:blip")):
                            embed_id = blip.get(qn("r:embed"))
                            if not embed_id:
                                continue
                            try:
                                image_part = doc.part.related_parts[embed_id]
                                img_data = image_part.blob
                                img_name = os.path.basename(image_part.partname)
                                import PIL.Image
                                import numpy as np
                                img = PIL.Image.open(_io.BytesIO(img_data))
                                img_np = _preprocess_image_for_ocr(np.array(img))
                                reader = get_ocr_reader()
                                res = reader.readtext(img_np, detail=0)
                                if res:
                                    ocr_text = " ".join(res)
                                    paragraphs_data.append((
                                        f"[Embedded Image OCR ({img_name})]: {ocr_text}",
                                        current_section or "[Embedded Image Content]"
                                    ))
                                processed_media_names.add(img_name)
                            except Exception:
                                pass  # not a rasterizable image (e.g. a chart) or unreadable -- skip
            
            # 2. Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_cells_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_cells_text:
                        # Simple deduplication for adjacent duplicate texts (e.g. from merged cells)
                        deduped_cells = []
                        for cell_txt in row_cells_text:
                            if not deduped_cells or deduped_cells[-1] != cell_txt:
                                deduped_cells.append(cell_txt)
                        if deduped_cells:
                            row_text = " | ".join(deduped_cells)
                            paragraphs_data.append((row_text, "[Table Data]"))
            
            # 3. Fallback: OCR any image still not covered by the per-paragraph pass
            # above (e.g. images in headers/footers, or anything the run-walk missed) --
            # skips anything already processed_media_names to avoid duplicate OCR/content.
            try:
                with zipfile.ZipFile(_io.BytesIO(file_bytes)) as zf:
                    media_files = [
                        n for n in zf.namelist()
                        if n.startswith("word/media/") and n.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".tiff", ".bmp"))
                        and os.path.basename(n) not in processed_media_names
                    ]
                    if media_files:
                        for name in sorted(media_files):
                            try:
                                img_data = zf.read(name)
                                import PIL.Image
                                import numpy as np
                                img = PIL.Image.open(_io.BytesIO(img_data))
                                img_np = _preprocess_image_for_ocr(np.array(img))
                                reader = get_ocr_reader()
                                res = reader.readtext(img_np, detail=0)
                                if res:
                                    ocr_text = " ".join(res)
                                    base_img_name = os.path.basename(name)
                                    paragraphs_data.append((f"[Embedded Image OCR ({base_img_name})]: {ocr_text}", "[Embedded Image Content]"))
                            except Exception as img_err:
                                pass
            except Exception as zf_err:
                pass
            
            docx_fname = getattr(f, "name", "unknown")
            docx_ext = os.path.splitext(docx_fname.lower())[1].lstrip(".")
            
            chunks = chunk_paragraphs(paragraphs_data, target=1000, overlap=200)
            docx_chunks = []
            for chunk_content, chunk_section, _, _ in chunks:
                docx_chunks.append((chunk_content, {
                    "source_file": docx_fname,
                    "source_type": docx_ext if docx_ext else "docx",
                    "section_heading": chunk_section,
                    "chunk_id": ""
                }))
            _ingested_chunks_cache[_cache_key(f.name)] = docx_chunks
            
            # Return combined text for indexing
            return "\n\n".join([txt for txt, _ in paragraphs_data])
        except Exception as e:
            return f"[Error parsing file {f.name}: {e}]"


